"""Text extraction from document files (PDF, DOCX, PPTX, TXT) and OCR for images."""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text(file_path: str, mime_type: str) -> str:
    """Extract text based on mime type."""
    if mime_type == "text/plain":
        return _extract_txt(file_path)
    elif mime_type == "application/pdf":
        return _extract_pdf(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx(file_path)
    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx(file_path)
    elif mime_type == "application/rtf":
        return _extract_rtf(file_path)
    else:
        logger.warning("Unsupported mime type for text extraction: %s", mime_type)
        return ""


def _extract_txt(file_path: str) -> str:
    return Path(file_path).read_text(errors="replace")


def _extract_pdf(file_path: str) -> str:
    try:
        import pymupdf
        doc = pymupdf.open(file_path)
        texts = []
        for page in doc:
            texts.append(page.get_text())
        doc.close()
        return "\n".join(texts)
    except ImportError:
        logger.warning("pymupdf not installed, cannot extract PDF text")
        return ""


def _extract_docx(file_path: str) -> str:
    """Extract text from .docx (and attempt .doc via textract fallback)."""
    ext = Path(file_path).suffix.lower()
    if ext == ".doc":
        return _extract_doc_legacy(file_path)
    try:
        from docx import Document
        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    texts.append(row_text)
        return "\n".join(texts)
    except ImportError:
        logger.warning("python-docx not installed, cannot extract DOCX text")
        return ""
    except Exception as e:
        logger.warning("Failed to extract DOCX text: %s", e)
        return ""


def _extract_doc_legacy(file_path: str) -> str:
    """Extract text from legacy .doc format using antiword or LibreOffice."""
    import subprocess
    # Try antiword first
    try:
        result = subprocess.run(["antiword", file_path], capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    # Try LibreOffice conversion
    try:
        import tempfile
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, file_path],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0:
                txt_file = Path(tmpdir) / (Path(file_path).stem + ".txt")
                if txt_file.exists():
                    return txt_file.read_text(errors="replace")
    except FileNotFoundError:
        pass
    logger.warning("Cannot extract legacy .doc — install antiword or LibreOffice")
    return ""


def _extract_pptx(file_path: str) -> str:
    """Extract text from .pptx (and attempt .ppt via LibreOffice fallback)."""
    ext = Path(file_path).suffix.lower()
    if ext == ".ppt":
        return _extract_legacy_office(file_path)
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPTX text")
        return ""
    except Exception as e:
        logger.warning("Failed to extract PPTX text: %s", e)
        return ""


def _extract_xlsx(file_path: str) -> str:
    """Extract text from .xlsx and .xls spreadsheets."""
    ext = Path(file_path).suffix.lower()
    if ext == ".xls":
        return _extract_legacy_office(file_path)
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            texts.append(f"--- Sheet: {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    texts.append(row_text)
        wb.close()
        return "\n".join(texts)
    except ImportError:
        logger.warning("openpyxl not installed, cannot extract XLSX text")
        return ""
    except Exception as e:
        logger.warning("Failed to extract XLSX text: %s", e)
        return ""


def _extract_rtf(file_path: str) -> str:
    """Extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text
        raw = Path(file_path).read_text(errors="replace")
        return rtf_to_text(raw)
    except ImportError:
        logger.warning("striprtf not installed, cannot extract RTF text")
        return ""
    except Exception as e:
        logger.warning("Failed to extract RTF text: %s", e)
        return ""


def _extract_legacy_office(file_path: str) -> str:
    """Fallback: convert legacy Office formats (.doc, .xls, .ppt) to text via LibreOffice."""
    import subprocess
    import tempfile
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, file_path],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0:
                txt_file = Path(tmpdir) / (Path(file_path).stem + ".txt")
                if txt_file.exists():
                    return txt_file.read_text(errors="replace")
    except FileNotFoundError:
        logger.warning("LibreOffice not available for legacy Office conversion")
    except Exception as e:
        logger.warning("Legacy Office conversion failed: %s", e)
    return ""


def ocr_image(file_path: str) -> str:
    """OCR an image file using pytesseract."""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(file_path)
        return pytesseract.image_to_string(img)
    except ImportError:
        logger.warning("pytesseract or Pillow not installed, cannot OCR image")
        return ""
